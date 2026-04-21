"""
One-time setup script: converts hardcoded Lead PPTX into a template
and copies the Ecommerce template into the templates/ folder.
Run once: python create_templates.py
"""
import shutil
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

LEAD_SOURCE = "/Users/mac0024/Downloads/05_11 - Report Semanal Emerald Turism Heath V3.pptx"
ECOMMERCE_SOURCE = "/Users/mac0024/Downloads/Report {{nome_cliente}} \u2013 {{periodo_analisado}}.pptx"
LEAD_OUT = "templates/template_lead.pptx"
ECOMMERCE_OUT = "templates/template_ecommerce.pptx"

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def find_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.shape_id == shape_id:
                    return child
    return None


def drop_extra_runs(para, keep_index=0):
    """Remove all runs after keep_index from paragraph XML."""
    runs = para.runs
    for run in runs[keep_index + 1:]:
        run._r.getparent().remove(run._r)


def set_single_run(shape, text, para_index=0):
    """Replace all runs in a paragraph with a single run containing text."""
    para = shape.text_frame.paragraphs[para_index]
    runs = para.runs
    if runs:
        runs[0].text = text
        drop_extra_runs(para, keep_index=0)
    else:
        # No runs exist — inject one via XML
        r_elem = etree.SubElement(para._p, f"{{{NS_A}}}r")
        rPr = etree.SubElement(r_elem, f"{{{NS_A}}}rPr")
        rPr.set("lang", "pt-BR")
        rPr.set("dirty", "0")
        t_elem = etree.SubElement(r_elem, f"{{{NS_A}}}t")
        t_elem.text = text


def inject_text(shape, text):
    """Add a run with text to an empty text box, copying rPr from sibling if possible."""
    para = shape.text_frame.paragraphs[0]
    r_elem = etree.SubElement(para._p, f"{{{NS_A}}}r")
    rPr = etree.SubElement(r_elem, f"{{{NS_A}}}rPr")
    rPr.set("lang", "pt-BR")
    rPr.set("dirty", "0")
    t_elem = etree.SubElement(r_elem, f"{{{NS_A}}}t")
    t_elem.text = text


def build_lead_template():
    prs = Presentation(LEAD_SOURCE)
    slide1 = prs.slides[0]
    slide2 = prs.slides[1]

    # ── SLIDE 1 ──────────────────────────────────────────────────
    # Shape 58: "Emerald Turism Heath" → {{nome_cliente}}
    set_single_run(find_shape(slide1, 58), "{{nome_cliente}}")

    # Shape 57: "Período: 05/12/2025 – 11/12/2025"
    # runs[0]="Período: " (keep bold), runs[1]="05", runs[2]="/12/2025 – 11/12/2025"
    s57 = find_shape(slide1, 57)
    para = s57.text_frame.paragraphs[0]
    runs = para.runs
    if len(runs) >= 2:
        runs[1].text = "{{periodo_analisado}}"
        drop_extra_runs(para, keep_index=1)

    # ── SLIDE 2 ──────────────────────────────────────────────────
    # Shape 92: "Período Analisado:\n05/12/2025 – 11/12/2025"
    s92 = find_shape(slide2, 92)
    para1 = s92.text_frame.paragraphs[1]
    runs = para1.runs
    if runs:
        runs[0].text = "{{periodo_analisado}}"
        drop_extra_runs(para1, keep_index=0)

    # Shape 86: "62" → {{lead_sem}}
    set_single_run(find_shape(slide2, 86), "{{lead_sem}}")

    # Shape 84 (empty, below lead): → "{{lead_sem_comp}} no período anterior"
    inject_text(find_shape(slide2, 84), "{{lead_sem_comp}} no período anterior")

    # Shape 88: "R$ 350,05" → {{inv_sem}}
    set_single_run(find_shape(slide2, 88), "{{inv_sem}}")

    # Shape 85 (empty, below inv): → "{{inv_sem_comp}} no período anterior"
    inject_text(find_shape(slide2, 85), "{{inv_sem_comp}} no período anterior")

    # Shape 98: "R$ 5,65" → {{cpl}}
    set_single_run(find_shape(slide2, 98), "{{cpl}}")

    # Shape 95 (empty, below cpl): → "{{cpl_comp}} no período anterior"
    inject_text(find_shape(slide2, 95), "{{cpl_comp}} no período anterior")

    # Groups — monthly section
    # Shape 74: "100" → {{leads_mes}}
    set_single_run(find_shape(slide2, 74), "{{leads_mes}}")
    # Shape 75: "- é a meta do mês" → {{meta_leads}} é a meta do mês
    set_single_run(find_shape(slide2, 75), "{{meta_leads}} é a meta do mês")
    # Shape 76: "- da meta foi atingida" → {{per_meta_leads}} da meta foi atingida
    set_single_run(find_shape(slide2, 76), "{{per_meta_leads}} da meta foi atingida")

    # Shape 80: "R$ 562,18" → {{inv_mes}}
    set_single_run(find_shape(slide2, 80), "{{inv_mes}}")
    # Shape 81: "- é a meta do mês" → {{meta_inv}} é a meta do mês
    set_single_run(find_shape(slide2, 81), "{{meta_inv}} é a meta do mês")
    # Shape 82: "- da meta foi atingida" → {{per_meta_inv}} da meta foi atingida
    set_single_run(find_shape(slide2, 82), "{{per_meta_inv}} da meta foi atingida")

    prs.save(LEAD_OUT)
    print(f"Saved: {LEAD_OUT}")


def main():
    build_lead_template()
    shutil.copy2(ECOMMERCE_SOURCE, ECOMMERCE_OUT)
    print(f"Copied: {ECOMMERCE_OUT}")
    print("Templates prontos!")


if __name__ == "__main__":
    main()
