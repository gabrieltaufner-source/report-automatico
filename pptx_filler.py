"""
Fills a PPTX template by replacing {{placeholders}} with real values.
Preserves all original formatting; colorizes comparison (▲/▼) text.
"""
import os
from datetime import datetime
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

GREEN = RGBColor(0x00, 0xB0, 0x50)
RED = RGBColor(0xFF, 0x00, 0x00)


def _color_for_text(text: str):
    """Return green for ▲, red for ▼, None otherwise."""
    if "▲" in text:
        return GREEN
    if "▼" in text:
        return RED
    return None


def _replace_in_para(para, data: dict) -> bool:
    """
    Replace any {{key}} found in a paragraph.
    Merges all runs into the first run after replacement to preserve style.
    Returns True if any replacement happened.
    """
    full_text = "".join(r.text for r in para.runs)
    changed = False

    for key, value in data.items():
        placeholder = "{{" + key + "}}"
        if placeholder in full_text:
            full_text = full_text.replace(placeholder, str(value))
            changed = True

    if not changed:
        return False

    runs = para.runs
    if not runs:
        return False

    # Keep first run's format, clear the rest
    first_run = runs[0]
    first_run.text = full_text
    for run in runs[1:]:
        run.text = ""

    # Apply ▲/▼ coloring
    color = _color_for_text(full_text)
    if color:
        first_run.font.color.rgb = color

    return True


def _process_shapes(shapes, data: dict):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _process_shapes(shape.shapes, data)
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            _replace_in_para(para, data)


def fill_template(tipo: str, data: dict, client_cfg: dict) -> str:
    """
    Open the correct template, fill all placeholders, and save.
    Returns the path of the saved file.
    """
    template_path = f"templates/template_{tipo}.pptx"
    if not os.path.exists(template_path):
        raise FileNotFoundError(
            f"Template não encontrado: {template_path}\n"
            "Execute 'python create_templates.py' primeiro."
        )

    prs = Presentation(template_path)

    for slide in prs.slides:
        _process_shapes(slide.shapes, data)

    # Build output file name: NomeCliente_DD-MM.pptx
    nome = client_cfg["nome"].replace(" ", "_")
    today = datetime.now().strftime("%d-%m")
    filename = f"{nome}_{today}.pptx"

    drive_folder = client_cfg.get("pasta_drive", "")
    if drive_folder and os.path.isdir(drive_folder):
        out_path = os.path.join(drive_folder, filename)
    else:
        os.makedirs("output", exist_ok=True)
        out_path = os.path.join("output", filename)

    prs.save(out_path)
    return out_path
